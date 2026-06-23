#!/usr/bin/env python3
"""
BSP Registry Tools -- PowerPoint Generator
==========================================
Generates a .pptx presentation from BSP Registry Tools content.

Dependencies:
    pip install python-pptx pillow matplotlib
    # for PPTX -> PDF conversion:
    # install libreoffice (soffice)

Usage:
    python3 generate_pptx.py
    # produces bsp_registry_tools.pptx in the same directory
    python3 generate_pptx.py --to-pdf
    # also produces bsp_registry_tools.pdf
    python3 generate_pptx.py --convert-only --input-pptx bsp_registry_tools.pptx
    # converts an existing PPTX to PDF
"""
from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
from pathlib import Path

HERE = Path(__file__).parent
ARCH_IMG_LEFT_INCHES = 0.7
ARCH_IMG_TOP_INCHES = 1.4
ARCH_IMG_WIDTH_INCHES = 8.6


def convert_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> None:
    if not pptx_path.exists():
        raise FileNotFoundError(f"Input PPTX not found: {pptx_path}")

    soffice_bin = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice_bin:
        raise RuntimeError(
            "LibreOffice is required for PPTX->PDF conversion. "
            "Install `soffice` (or `libreoffice`) and retry."
        )

    pptx_path = pptx_path.resolve()
    outdir = pdf_path.parent.resolve()
    outdir.mkdir(parents=True, exist_ok=True)

    cmd = [
        soffice_bin,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(outdir),
        str(pptx_path),
    ]
    subprocess.run(cmd, check=True)

    converted_pdf = outdir / f"{pptx_path.stem}.pdf"
    if converted_pdf != pdf_path.resolve() and converted_pdf.exists():
        converted_pdf.replace(pdf_path.resolve())
    if not pdf_path.exists():
        raise RuntimeError(f"PPTX->PDF conversion did not produce: {pdf_path}")

    print(f"  Generated {pdf_path} ({pdf_path.stat().st_size:,} bytes)")


def build_pptx(output: Path, arch_img: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Layout aliases
    title_layout = prs.slide_layouts[0]      # Title Slide
    title_body_layout = prs.slide_layouts[1] # Title and Content

    def add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(title_body_layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for idx, item in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(22)

    add_title_slide(
        "BSP Registry Tools",
        "Streamlined Yocto and Isar BSP management\nfor embedded Linux teams",
    )

    add_bullets_slide(
        "What is BSP Registry Tools?",
        [
            "Unified CLI and API for BSP build workflows",
            "YAML registry as a single source of truth",
            "Supports Yocto and Isar release flows",
            "Designed for reproducibility and team collaboration",
        ],
    )

    add_bullets_slide(
        "Core Benefits",
        [
            "Faster onboarding for new engineers and projects",
            "Repeatable builds across devices and releases",
            "Reduced manual errors from ad-hoc build scripts",
            "Improved traceability from config to artifact",
        ],
    )

    add_bullets_slide(
        "Latest Improvements",
        [
            "Build provenance: build-manifest.json generated per build",
            "Repo-manifest export for downstream integration flows",
            "Runtime selection flags with shell completion support",
            "Expanded testing backends and improved direct-test flows",
            "Flexible scan and SBOM workflows for CI pipelines",
            "Multi-registry and named remote collaboration",
        ],
    )

    # Architecture slide (optional image)
    slide = prs.slides.add_slide(title_body_layout)
    slide.shapes.title.text = "System Architecture"
    if arch_img.exists():
        slide.shapes.add_picture(
            str(arch_img),
            Inches(ARCH_IMG_LEFT_INCHES),
            Inches(ARCH_IMG_TOP_INCHES),
            width=Inches(ARCH_IMG_WIDTH_INCHES),
        )
    else:
        slide.placeholders[1].text = "Architecture image not available"

    add_bullets_slide(
        "Typical Workflow",
        [
            "Select preset (or device + release + features)",
            "Build via KAS/kas-container orchestration",
            "Deploy, gather, scan, and test artifacts",
            "Integrate with REST/GraphQL APIs and CI automation",
        ],
    )

    add_bullets_slide(
        "Adoption Path",
        [
            "Start with one product line",
            "Define shared presets and release variants",
            "Roll into CI/CD build, test, and scan pipelines",
            "Expand to multi-team multi-registry operations",
        ],
    )

    add_bullets_slide(
        "Summary",
        [
            "Accelerates delivery while improving reliability",
            "Aligns engineering workflows across teams",
            "Scales from local developer use to enterprise CI/CD",
        ],
    )

    prs.save(str(output))
    print(f"  Generated {output} ({output.stat().st_size:,} bytes, {len(prs.slides)} slides)")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate BSP Registry Tools PPTX and optional PDF.")
    parser.add_argument("--to-pdf", action="store_true", help="Also convert generated PPTX to PDF.")
    parser.add_argument(
        "--convert-only",
        action="store_true",
        help="Only convert an existing PPTX to PDF (skip PPTX generation).",
    )
    parser.add_argument(
        "--input-pptx",
        type=Path,
        default=HERE / "bsp_registry_tools.pptx",
        help="Input PPTX path for conversion (default: presentation/bsp_registry_tools.pptx).",
    )
    parser.add_argument(
        "--pdf-output",
        type=Path,
        default=HERE / "bsp_registry_tools.pdf",
        help="PDF output path for conversion (default: presentation/bsp_registry_tools.pdf).",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()

    if args.convert_only:
        print("[1/1] Converting PPTX to PDF …")
        convert_pptx_to_pdf(args.input_pptx, args.pdf_output)
        print("Done.")
        return

    # Reuse the same architecture diagram generation as PDF tooling.
    sys.path.insert(0, str(HERE))
    from generate_pdf import generate_architecture_diagram

    arch_img = HERE / "images" / "architecture.png"
    output = HERE / "bsp_registry_tools.pptx"

    print("[1/2] Generating architecture diagram …")
    generate_architecture_diagram(arch_img)

    print("[2/2] Building PPTX presentation …")
    build_pptx(output, arch_img)

    if args.to_pdf:
        print("[3/3] Converting PPTX to PDF …")
        convert_pptx_to_pdf(output, args.pdf_output)

    print("Done.")


if __name__ == "__main__":
    main()
